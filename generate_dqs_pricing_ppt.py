"""
InvenIQ — Design Quote Studio Subscription & Licensing
Generates a professional 12-slide PowerPoint presentation.
Output: c:\InvenIQ\DQS_Subscription_Pricing.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Colour palette ─────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0d, 0x1a, 0x14)   # dark green-black
BG_CARD   = RGBColor(0x11, 0x1f, 0x18)   # card surface
BG_LIGHT  = RGBColor(0xf8, 0xfa, 0xf8)   # slide bg (light slides)
WHITE     = RGBColor(0xff, 0xff, 0xff)
GREEN     = RGBColor(0x10, 0xb9, 0x81)   # InvenIQ green
GREEN_DK  = RGBColor(0x05, 0x7a, 0x55)
TEAL      = RGBColor(0x0f, 0x76, 0x6e)   # Starter
BLUE      = RGBColor(0x1d, 0x4e, 0xd8)   # Professional
PURPLE    = RGBColor(0x7c, 0x3a, 0xed)   # Enterprise
AMBER     = RGBColor(0xd9, 0x77, 0x06)
RED       = RGBColor(0xdc, 0x26, 0x26)
TEXT_MID  = RGBColor(0x8a, 0xab, 0x91)
TEXT_DIM  = RGBColor(0x5c, 0x7a, 0x62)
BORDER    = RGBColor(0x1d, 0x32, 0x25)
SLATE     = RGBColor(0x1e, 0x29, 0x3b)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, border_rgb=None, border_width=Pt(0)):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.line.width = border_width
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if border_rgb:
        shp.line.color.rgb = border_rgb
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, x, y, w, h, size=Pt(12), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True, font_name='Calibri'):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb


def add_label_value(slide, label, value, x, y, w,
                    label_color=TEXT_MID, value_color=WHITE,
                    label_size=Pt(9), value_size=Pt(18), font_name='Calibri'):
    add_text(slide, label, x, y, w, Inches(.22), size=label_size,
             color=label_color, font_name=font_name)
    add_text(slide, value, x, y + Inches(.23), w, Inches(.36), size=value_size,
             bold=True, color=value_color, font_name=font_name)


def slide_bg(slide, color=BG_DARK):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=color)


def top_bar(slide, accent=GREEN):
    add_rect(slide, 0, 0, SLIDE_W, Inches(.07), fill_rgb=accent)


def section_tag(slide, text, x, y, color=GREEN):
    r = add_rect(slide, x, y, Inches(2.1), Inches(.26), fill_rgb=None)
    r.fill.background()
    r.line.fill.background()
    add_text(slide, text, x, y, Inches(2.1), Inches(.26),
             size=Pt(8), bold=True, color=color, font_name='Courier New')


def card(slide, x, y, w, h, accent=GREEN, dark=True):
    bg = BG_CARD if dark else RGBColor(0xf0, 0xfd, 0xf4)
    shp = add_rect(slide, x, y, w, h, fill_rgb=bg, border_rgb=BORDER, border_width=Pt(0.75))
    add_rect(slide, x, y, w, Inches(.04), fill_rgb=accent)
    return shp


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)

# green left panel
add_rect(sl, 0, 0, Inches(4.8), SLIDE_H, fill_rgb=RGBColor(0x05, 0x2e, 0x16))
add_rect(sl, 0, 0, Inches(.06), SLIDE_H, fill_rgb=GREEN)

# logo / brand text
add_text(sl, 'InvenIQ', Inches(.22), Inches(.5), Inches(4.2), Inches(.55),
         size=Pt(28), bold=True, color=GREEN, font_name='Calibri')
add_text(sl, 'AI-Powered Inventory Intelligence Platform',
         Inches(.22), Inches(1.05), Inches(4.2), Inches(.3),
         size=Pt(9), color=TEXT_MID, font_name='Calibri')

# divider
add_rect(sl, Inches(.22), Inches(1.45), Inches(4.2), Pt(1), fill_rgb=BORDER)

# main title
add_text(sl, 'Design Quote Studio',
         Inches(.22), Inches(1.65), Inches(4.5), Inches(.65),
         size=Pt(30), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'Subscription & Licensing',
         Inches(.22), Inches(2.3), Inches(4.5), Inches(.45),
         size=Pt(20), bold=False, color=GREEN, font_name='Calibri')

# subtitle bullets
for i, line in enumerate([
    'Per-month pricing for DQS module + AI components',
    'Build-in licensing cost transparency',
    'Interior Quotations · Architect Fee Proposals · AI BOQ Scanner',
]):
    add_text(sl, f'→  {line}',
             Inches(.22), Inches(2.95) + Inches(.38) * i, Inches(4.4), Inches(.36),
             size=Pt(10), color=TEXT_MID, font_name='Calibri')

# version / date
add_text(sl, 'v3.7  ·  June 2026  ·  Architect Role Edition',
         Inches(.22), Inches(6.8), Inches(4.5), Inches(.3),
         size=Pt(8.5), color=TEXT_DIM, font_name='Courier New')

# right panel — big metric tiles
metrics = [
    ('3',           'Subscription Tiers',    GREEN),
    ('₹2,999',      'Starting / month',      TEAL),
    ('₹5,999',      'Professional / month',  BLUE),
    ('₹11,999',     'Enterprise / month',    PURPLE),
]
for idx, (val, lbl, col) in enumerate(metrics):
    tx = Inches(5.3)
    ty = Inches(.55) + Inches(1.55) * idx
    tw = Inches(7.5)
    th = Inches(1.35)
    card(sl, tx, ty, tw, th, accent=col)
    add_text(sl, val, tx + Inches(.22), ty + Inches(.2), tw - Inches(.4), Inches(.6),
             size=Pt(36), bold=True, color=col, font_name='Calibri')
    add_text(sl, lbl, tx + Inches(.22), ty + Inches(.78), tw - Inches(.4), Inches(.3),
             size=Pt(11), color=TEXT_MID, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS DESIGN QUOTE STUDIO?
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, GREEN)

add_text(sl, 'What is Design Quote Studio?',
         Inches(.45), Inches(.3), Inches(9), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'InvenIQ  ·  Module: designquote  ·  Role: Architect (restricted)',
         Inches(.45), Inches(.82), Inches(9), Inches(.28),
         size=Pt(9), color=TEXT_MID, font_name='Courier New')

caps = [
    ('🏛',  'Interior Quotations',        'Create detailed room-wise BOQ quotations for interior design projects.\nDraft, Sent, Negotiating, Won, Lost, Expired workflow stages.',  TEAL),
    ('📐',  'Architect Fee Proposals',     'Generate structured fee proposals (5–8% of project cost) across\n6 phases: P1 Concept → P6 Post-Handover. Auto GST at 18% (SAC 998331).',  BLUE),
    ('🤖',  'AI BOQ Scanner',             'WhatsApp photo → AI identifies CP fittings, sanitary ware, hardware,\ntiles, plumbing items with HSN codes and quantities instantly.',  GREEN),
    ('🔗',  'QB↔DQS Two-Way Sync',        'Quotation Builder items sync into Design Quote Studio BOQ sections.\nProducts saved to shared Product Catalog with AI-inferred HSN + category.',  PURPLE),
    ('📄',  'PDF Export & Proposals',      'Professional PDF output for client delivery. Custom letterhead\nand branding on Enterprise plan. No per-PDF API cost.',  AMBER),
    ('🗂',  'Product Catalog Integration', 'All scanned items auto-saved to the shared catalog.\nBrand, finish, HSN, and category inferred by AI on every scan.',  RGBColor(0x06, 0x78, 0x60)),
]

cols = 3
for i, (icon, title, desc, col) in enumerate(caps):
    row = i // cols
    column = i % cols
    cx = Inches(.38) + Inches(4.27) * column
    cy = Inches(1.32) + Inches(2.5) * row
    cw = Inches(4.08)
    ch = Inches(2.3)
    card(sl, cx, cy, cw, ch, accent=col)
    add_text(sl, icon + '  ' + title,
             cx + Inches(.18), cy + Inches(.22), cw - Inches(.3), Inches(.35),
             size=Pt(13), bold=True, color=col, font_name='Calibri')
    add_text(sl, desc,
             cx + Inches(.18), cy + Inches(.62), cw - Inches(.3), Inches(1.5),
             size=Pt(9.5), color=TEXT_MID, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SUBSCRIPTION TIERS OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, GREEN)

add_text(sl, 'Subscription Tiers — Monthly Pricing',
         Inches(.45), Inches(.3), Inches(10), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'Three plans · All billed monthly · Cancel anytime · 30-day free trial on Professional',
         Inches(.45), Inches(.82), Inches(10), Inches(.28),
         size=Pt(9), color=TEXT_MID, font_name='Calibri')

tiers = [
    {
        'name': 'STARTER',
        'price': '₹2,999',
        'seats': '1 Architect seat',
        'color': TEAL,
        'features': [
            '30 Interior Quotations / month',
            '5 Architect Fee Proposals / month',
            'AI BOQ Scanner — 50 scans / month',
            'WhatsApp photo scan — 5 / month',
            'Standard PDF export',
            '✗  Product Catalog sync',
            '✗  QB↔DQS two-way sync',
            '✗  Custom PDF branding',
        ],
    },
    {
        'name': 'PROFESSIONAL',
        'price': '₹5,999',
        'seats': '3 Architect seats',
        'color': BLUE,
        'badge': '★  RECOMMENDED',
        'features': [
            'Unlimited Interior Quotations',
            'Unlimited Architect Fee Proposals',
            'AI BOQ Scanner — unlimited scans',
            'WhatsApp + multi-photo — unlimited',
            'Standard PDF export',
            'Product Catalog sync (hardware/sanitary)',
            'QB↔DQS two-way sync',
            '✗  Custom PDF branding',
        ],
    },
    {
        'name': 'ENTERPRISE',
        'price': '₹11,999',
        'seats': '10 Architect seats',
        'color': PURPLE,
        'features': [
            'Unlimited Interior Quotations',
            'Unlimited Architect Fee Proposals',
            'AI BOQ Scanner — unlimited scans',
            'WhatsApp + multi-photo — unlimited',
            'Custom PDF branding & letterhead',
            'Product Catalog + QB↔DQS sync',
            '2 custom template designs (one-time)',
            'Dedicated account manager + 99.9% SLA',
        ],
    },
]

for idx, t in enumerate(tiers):
    cx = Inches(.38) + Inches(4.28) * idx
    cy = Inches(1.18)
    cw = Inches(4.08)
    ch = Inches(5.95)
    col = t['color']

    card(sl, cx, cy, cw, ch, accent=col)

    # recommended badge
    if 'badge' in t:
        add_rect(sl, cx + Inches(1.0), cy - Inches(.18), Inches(2.1), Inches(.27),
                 fill_rgb=BLUE)
        add_text(sl, t['badge'],
                 cx + Inches(1.0), cy - Inches(.18), Inches(2.1), Inches(.27),
                 size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 font_name='Courier New')

    add_text(sl, t['name'],
             cx + Inches(.2), cy + Inches(.18), cw - Inches(.4), Inches(.28),
             size=Pt(10), bold=True, color=col, font_name='Courier New')
    add_text(sl, t['price'],
             cx + Inches(.2), cy + Inches(.48), cw - Inches(.4), Inches(.55),
             size=Pt(30), bold=True, color=WHITE, font_name='Calibri')
    add_text(sl, '/ month',
             cx + Inches(.2), cy + Inches(1.02), cw - Inches(.4), Inches(.25),
             size=Pt(9), color=TEXT_MID, font_name='Calibri')
    add_text(sl, t['seats'],
             cx + Inches(.2), cy + Inches(1.26), cw - Inches(.4), Inches(.25),
             size=Pt(9.5), bold=True, color=col, font_name='Calibri')

    add_rect(sl, cx + Inches(.2), cy + Inches(1.6), cw - Inches(.4), Pt(0.8),
             fill_rgb=BORDER)

    for fi, feat in enumerate(t['features']):
        fc = TEXT_DIM if feat.startswith('✗') else TEXT_MID
        fy = cy + Inches(1.72) + Inches(.35) * fi
        add_text(sl, ('→  ' if not feat.startswith('✗') else '      ') + feat,
                 cx + Inches(.18), fy, cw - Inches(.32), Inches(.34),
                 size=Pt(9), color=fc, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STARTER PLAN DETAIL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, TEAL)

add_rect(sl, 0, 0, Inches(4.0), SLIDE_H, fill_rgb=RGBColor(0x07, 0x2a, 0x28))
add_rect(sl, 0, 0, Inches(.06), SLIDE_H, fill_rgb=TEAL)

add_text(sl, 'STARTER', Inches(.22), Inches(.4), Inches(3.5), Inches(.35),
         size=Pt(11), bold=True, color=TEAL, font_name='Courier New')
add_text(sl, '₹2,999', Inches(.22), Inches(.78), Inches(3.5), Inches(.65),
         size=Pt(42), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'per month', Inches(.22), Inches(1.42), Inches(3.5), Inches(.28),
         size=Pt(11), color=TEXT_MID, font_name='Calibri')

add_rect(sl, Inches(.22), Inches(1.82), Inches(3.5), Pt(0.8), fill_rgb=BORDER)

specs = [
    ('Architect Seats',          '1 seat included'),
    ('Interior Quotations',      '30 / month'),
    ('Architect Fee Proposals',  '5 / month'),
    ('AI BOQ Scans',             '50 / month'),
    ('WhatsApp Photo Scans',     '5 / month'),
    ('PDF Export',               'Standard template'),
    ('Catalog Sync',             '✗  Not included'),
    ('QB↔DQS Sync',              '✗  Not included'),
]
for i, (k, v) in enumerate(specs):
    ty = Inches(2.0) + Inches(.5) * i
    add_text(sl, k, Inches(.22), ty, Inches(2.1), Inches(.4),
             size=Pt(10), color=TEXT_DIM, font_name='Calibri')
    vc = TEXT_DIM if v.startswith('✗') else WHITE
    add_text(sl, v, Inches(2.35), ty, Inches(1.5), Inches(.4),
             size=Pt(10), bold=True, color=vc, font_name='Calibri')

# right panel — best for section
add_text(sl, 'Best for:', Inches(4.3), Inches(.38), Inches(8.7), Inches(.3),
         size=Pt(11), bold=True, color=TEAL, font_name='Calibri')
add_text(sl, 'Independent architects and freelance interior designers handling\n'
             'low-to-medium volume projects. Ideal for trying out InvenIQ DQS\n'
             'before committing to a higher tier.',
         Inches(4.3), Inches(.72), Inches(8.7), Inches(.95),
         size=Pt(11), color=TEXT_MID, font_name='Calibri')

use_cases = [
    ('Single architect firm',     'Perfect for a solo practitioner managing 20–30 projects/year'),
    ('New to digital quotations', 'Replaces Excel/WhatsApp-based quotation workflows'),
    ('Low proposal volume',       'Up to 5 detailed fee proposals per month'),
    ('Trial before upgrade',      'Full platform access to evaluate before moving to Professional'),
]
for i, (title, desc) in enumerate(use_cases):
    uy = Inches(1.85) + Inches(1.35) * i
    card(sl, Inches(4.3), uy, Inches(8.7), Inches(1.18), accent=TEAL)
    add_text(sl, title, Inches(4.52), uy + Inches(.18), Inches(8.3), Inches(.32),
             size=Pt(12), bold=True, color=TEAL, font_name='Calibri')
    add_text(sl, desc, Inches(4.52), uy + Inches(.5), Inches(8.3), Inches(.55),
             size=Pt(10.5), color=TEXT_MID, font_name='Calibri')

add_text(sl, '→  Upgrade to Professional when scan volume or proposal count exceeds limits',
         Inches(4.3), Inches(7.12), Inches(8.7), Inches(.28),
         size=Pt(9), color=TEXT_DIM, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROFESSIONAL PLAN DETAIL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, BLUE)

add_rect(sl, 0, 0, Inches(4.0), SLIDE_H, fill_rgb=RGBColor(0x0c, 0x1a, 0x3a))
add_rect(sl, 0, 0, Inches(.06), SLIDE_H, fill_rgb=BLUE)

add_text(sl, '★  PROFESSIONAL  ·  RECOMMENDED', Inches(.22), Inches(.38), Inches(3.7), Inches(.3),
         size=Pt(9), bold=True, color=BLUE, font_name='Courier New')
add_text(sl, '₹5,999', Inches(.22), Inches(.72), Inches(3.5), Inches(.65),
         size=Pt(42), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'per month', Inches(.22), Inches(1.36), Inches(3.5), Inches(.28),
         size=Pt(11), color=TEXT_MID, font_name='Calibri')
add_rect(sl, Inches(.22), Inches(1.76), Inches(3.5), Pt(0.8), fill_rgb=BORDER)

specs = [
    ('Architect Seats',          '3 seats'),
    ('Interior Quotations',      'Unlimited'),
    ('Architect Fee Proposals',  'Unlimited'),
    ('AI BOQ Scans',             'Unlimited'),
    ('WhatsApp Photo Scans',     'Unlimited'),
    ('PDF Export',               'Standard template'),
    ('Catalog Sync',             '✓  Included'),
    ('QB↔DQS Sync',              '✓  Included'),
]
for i, (k, v) in enumerate(specs):
    ty = Inches(1.92) + Inches(.5) * i
    add_text(sl, k, Inches(.22), ty, Inches(2.1), Inches(.4),
             size=Pt(10), color=TEXT_DIM, font_name='Calibri')
    vc = GREEN if v.startswith('✓') else (TEXT_DIM if v.startswith('✗') else WHITE)
    add_text(sl, v, Inches(2.35), ty, Inches(1.5), Inches(.4),
             size=Pt(10), bold=True, color=vc, font_name='Calibri')

add_text(sl, 'Best for:', Inches(4.3), Inches(.38), Inches(8.7), Inches(.3),
         size=Pt(11), bold=True, color=BLUE, font_name='Calibri')
add_text(sl, 'Growing interior design firms and architecture studios managing\n'
             'multiple projects simultaneously with a small team of architects.',
         Inches(4.3), Inches(.72), Inches(8.7), Inches(.7),
         size=Pt(11), color=TEXT_MID, font_name='Calibri')

use_cases = [
    ('Multi-architect studio',      '2–3 architects sharing the same DQS workspace and catalog'),
    ('High quotation volume',        'No cap on Interior Quotations or Fee Proposals per month'),
    ('Catalog-driven workflow',      'Product Catalog synced with QB scanner — consistent pricing'),
    ('Cross-module efficiency',      'QB↔DQS sync eliminates double-entry of materials data'),
]
for i, (title, desc) in enumerate(use_cases):
    uy = Inches(1.62) + Inches(1.35) * i
    card(sl, Inches(4.3), uy, Inches(8.7), Inches(1.18), accent=BLUE)
    add_text(sl, title, Inches(4.52), uy + Inches(.18), Inches(8.3), Inches(.32),
             size=Pt(12), bold=True, color=BLUE, font_name='Calibri')
    add_text(sl, desc, Inches(4.52), uy + Inches(.5), Inches(8.3), Inches(.55),
             size=Pt(10.5), color=TEXT_MID, font_name='Calibri')

add_text(sl, '→  Most popular plan · Includes 30-day free trial for all new installations',
         Inches(4.3), Inches(7.12), Inches(8.7), Inches(.28),
         size=Pt(9), color=BLUE, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ENTERPRISE PLAN DETAIL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, PURPLE)

add_rect(sl, 0, 0, Inches(4.0), SLIDE_H, fill_rgb=RGBColor(0x18, 0x0a, 0x35))
add_rect(sl, 0, 0, Inches(.06), SLIDE_H, fill_rgb=PURPLE)

add_text(sl, 'ENTERPRISE', Inches(.22), Inches(.38), Inches(3.7), Inches(.3),
         size=Pt(11), bold=True, color=PURPLE, font_name='Courier New')
add_text(sl, '₹11,999', Inches(.22), Inches(.72), Inches(3.5), Inches(.65),
         size=Pt(42), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'per month', Inches(.22), Inches(1.36), Inches(3.5), Inches(.28),
         size=Pt(11), color=TEXT_MID, font_name='Calibri')
add_rect(sl, Inches(.22), Inches(1.76), Inches(3.5), Pt(0.8), fill_rgb=BORDER)

specs = [
    ('Architect Seats',           '10 seats included'),
    ('Interior Quotations',       'Unlimited'),
    ('Architect Fee Proposals',   'Unlimited'),
    ('AI BOQ Scans',              'Unlimited'),
    ('WhatsApp Photo Scans',      'Unlimited'),
    ('Custom PDF Branding',       '✓  Included'),
    ('Template Design',           '2 templates (one-time)'),
    ('Account Manager',           '✓  Dedicated'),
    ('SLA',                       '99.9% uptime'),
]
for i, (k, v) in enumerate(specs):
    ty = Inches(1.88) + Inches(.47) * i
    add_text(sl, k, Inches(.22), ty, Inches(2.1), Inches(.38),
             size=Pt(9.5), color=TEXT_DIM, font_name='Calibri')
    vc = GREEN if v.startswith('✓') else WHITE
    add_text(sl, v, Inches(2.35), ty, Inches(1.55), Inches(.38),
             size=Pt(9.5), bold=True, color=vc, font_name='Calibri')

add_text(sl, 'Best for:', Inches(4.3), Inches(.38), Inches(8.7), Inches(.3),
         size=Pt(11), bold=True, color=PURPLE, font_name='Calibri')
add_text(sl, 'Large interior design firms, real-estate developers, and institutional buyers\n'
             'running multiple projects across teams with brand consistency requirements.',
         Inches(4.3), Inches(.72), Inches(8.7), Inches(.75),
         size=Pt(11), color=TEXT_MID, font_name='Calibri')

use_cases = [
    ('Large team deployment',       'Up to 10 architect logins with role-scoped module access'),
    ('Brand identity in proposals', 'Custom letterhead, logo, and colour scheme on all PDFs'),
    ('Template ownership',          '2 bespoke templates designed by InvenIQ team at no extra cost'),
    ('Dedicated support',           'Named account manager · Priority SLA · 99.9% uptime guarantee'),
]
for i, (title, desc) in enumerate(use_cases):
    uy = Inches(1.62) + Inches(1.35) * i
    card(sl, Inches(4.3), uy, Inches(8.7), Inches(1.18), accent=PURPLE)
    add_text(sl, title, Inches(4.52), uy + Inches(.18), Inches(8.3), Inches(.32),
             size=Pt(12), bold=True, color=PURPLE, font_name='Calibri')
    add_text(sl, desc, Inches(4.52), uy + Inches(.5), Inches(8.3), Inches(.55),
             size=Pt(10.5), color=TEXT_MID, font_name='Calibri')

add_text(sl, '→  Additional seats: ₹999/seat/month beyond the 10 included',
         Inches(4.3), Inches(7.12), Inches(8.7), Inches(.28),
         size=Pt(9), color=PURPLE, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURE COMPARISON TABLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, GREEN)

add_text(sl, 'Feature Comparison — All Tiers',
         Inches(.45), Inches(.3), Inches(12), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')

headers = ['Feature', 'Starter\n₹2,999/mo', 'Professional\n₹5,999/mo', 'Enterprise\n₹11,999/mo']
cols_x  = [Inches(.38), Inches(4.45), Inches(7.4), Inches(10.35)]
col_w   = [Inches(3.9),  Inches(2.82), Inches(2.82), Inches(2.82)]
hdr_colors = [BG_CARD, TEAL, BLUE, PURPLE]

rows = [
    ('Architect Seats',           '1',         '3',           '10'),
    ('Interior Quotations/month', '30',        'Unlimited',   'Unlimited'),
    ('Fee Proposals/month',       '5',         'Unlimited',   'Unlimited'),
    ('AI BOQ Scans/month',        '50',        'Unlimited',   'Unlimited'),
    ('WhatsApp Photo Scans',      '5/month',   'Unlimited',   'Unlimited'),
    ('PDF Export',                '✓',         '✓',           '✓'),
    ('Custom PDF Branding',       '✗',         '✗',           '✓'),
    ('Product Catalog Sync',      '✗',         '✓',           '✓'),
    ('QB↔DQS Sync',               '✗',         '✓',           '✓'),
    ('Custom Templates',          '✗',         '✗',           '2 (one-time)'),
    ('Account Manager',           '✗',         '✗',           '✓ Dedicated'),
    ('Uptime SLA',                'Standard',  'Standard',    '99.9% SLA'),
]

ROW_H = Inches(.43)
HDR_H = Inches(.56)
TBL_Y = Inches(1.0)

# header row
for ci, (hdr, cx, cw, hcol) in enumerate(zip(headers, cols_x, col_w, hdr_colors)):
    add_rect(sl, cx, TBL_Y, cw, HDR_H, fill_rgb=hcol, border_rgb=BORDER, border_width=Pt(0.5))
    align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
    padx = Inches(.12) if ci > 0 else Inches(.14)
    add_text(sl, hdr, cx + padx, TBL_Y + Inches(.05), cw - padx, HDR_H - Inches(.08),
             size=Pt(9.5), bold=True, color=WHITE, align=align, font_name='Calibri')

for ri, row in enumerate(rows):
    ry = TBL_Y + HDR_H + ROW_H * ri
    bg = RGBColor(0x0f, 0x22, 0x18) if ri % 2 == 0 else BG_CARD
    for ci, (cell, cx, cw) in enumerate(zip(row, cols_x, col_w)):
        add_rect(sl, cx, ry, cw, ROW_H, fill_rgb=bg, border_rgb=BORDER, border_width=Pt(0.3))
        align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        padx = Inches(.1) if ci > 0 else Inches(.14)
        cc = GREEN if cell == '✓' or cell.startswith('✓') else (TEXT_DIM if cell == '✗' else (WHITE if ci == 0 else TEXT_MID))
        add_text(sl, cell, cx + padx, ry + Inches(.06), cw - padx, ROW_H - Inches(.1),
                 size=Pt(9), bold=(ci == 0), color=cc, align=align, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ADD-ON COMPONENTS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, AMBER)

add_text(sl, 'Add-on Components',
         Inches(.45), Inches(.3), Inches(10), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'Optional modules available with any subscription tier',
         Inches(.45), Inches(.82), Inches(10), Inches(.28),
         size=Pt(9), color=TEXT_MID, font_name='Calibri')

addons = [
    {
        'icon': '👤',
        'title': 'Extra Architect Seat',
        'price': '₹799',
        'unit': '/ seat / month',
        'note': 'Add architect user accounts beyond the included seats in your plan.\nEach seat gets full DQS access under your organisation.',
        'for': 'Starter · Professional · Enterprise',
        'col': TEAL,
    },
    {
        'icon': '🤖',
        'title': 'AI Scan Pack — 500 extra scans',
        'price': '₹1,499',
        'unit': '/ month',
        'note': 'Purchase additional 500 AI BOQ scans per month on top of your plan limit.\nOnly applicable on the Starter plan (Professional and Enterprise are unlimited).',
        'for': 'Starter only',
        'col': GREEN,
    },
    {
        'icon': '📄',
        'title': 'Custom PDF Template Design',
        'price': '₹2,999',
        'unit': 'one-time',
        'note': 'InvenIQ design team creates a bespoke quotation PDF template\nmatching your brand: logo, colours, fonts, header/footer layout.',
        'for': 'Starter · Professional',
        'col': BLUE,
    },
    {
        'icon': '💬',
        'title': 'WhatsApp Bot Integration',
        'price': '₹1,999',
        'unit': '/ month',
        'note': 'Dedicated WhatsApp Business API integration. Clients send site photos\ndirectly to your bot — AI auto-generates BOQ and saves to DQS.',
        'for': 'All tiers',
        'col': PURPLE,
    },
]

for i, a in enumerate(addons):
    row = i // 2
    col = i % 2
    ax = Inches(.38) + Inches(6.47) * col
    ay = Inches(1.28) + Inches(2.9) * row
    aw = Inches(6.28)
    ah = Inches(2.65)
    card(sl, ax, ay, aw, ah, accent=a['col'])
    add_text(sl, a['icon'] + '  ' + a['title'],
             ax + Inches(.2), ay + Inches(.2), aw - Inches(.4), Inches(.38),
             size=Pt(14), bold=True, color=a['col'], font_name='Calibri')
    # price block
    add_text(sl, a['price'],
             ax + Inches(.2), ay + Inches(.62), Inches(1.6), Inches(.5),
             size=Pt(28), bold=True, color=WHITE, font_name='Calibri')
    add_text(sl, a['unit'],
             ax + Inches(1.82), ay + Inches(.78), Inches(1.4), Inches(.32),
             size=Pt(9.5), color=TEXT_MID, font_name='Calibri')
    add_text(sl, a['note'],
             ax + Inches(.2), ay + Inches(1.18), aw - Inches(.38), Inches(.85),
             size=Pt(9.5), color=TEXT_MID, font_name='Calibri')
    add_text(sl, 'Available on: ' + a['for'],
             ax + Inches(.2), ay + Inches(2.25), aw - Inches(.38), Inches(.28),
             size=Pt(8.5), color=a['col'], font_name='Courier New')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUILD-IN LICENSING COST REFERENCE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, AMBER)

add_text(sl, 'Build-in Licensing Cost Reference',
         Inches(.45), Inches(.3), Inches(10), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'What is included in the InvenIQ DQS build · What you pay separately to OpenAI at cost',
         Inches(.45), Inches(.82), Inches(10), Inches(.28),
         size=Pt(9), color=TEXT_MID, font_name='Calibri')

costs = [
    {
        'label': 'DQS Core Engine',
        'cost': 'INCLUDED',
        'note': 'Fully licensed as part of InvenIQ Enterprise.\nNo separate module activation fee.',
        'col': GREEN, 'included': True,
    },
    {
        'label': 'PDF Generation',
        'cost': 'INCLUDED',
        'note': 'Client-side rendering (browser-based).\nZero server cost per PDF.',
        'col': GREEN, 'included': True,
    },
    {
        'label': 'Quote & Proposal Storage',
        'cost': 'INCLUDED',
        'note': 'Stored in your own MySQL database.\nNo cloud storage fee from InvenIQ.',
        'col': GREEN, 'included': True,
    },
    {
        'label': 'GPT-4o API (per AI scan)',
        'cost': '~₹8–12 / call',
        'note': '$0.10–$0.15 per scan · Billed directly\nto your OpenAI account at cost.',
        'col': AMBER, 'included': False,
    },
    {
        'label': 'WhatsApp Photo Parsing',
        'cost': '~₹2–4 / image',
        'note': 'Vision model tokens · Charged per image.\nBilled to your OpenAI account.',
        'col': AMBER, 'included': False,
    },
    {
        'label': 'Catalog AI Inference',
        'cost': '~₹1–2 / lookup',
        'note': 'HSN + category per product match.\nGPT-4o mini — very low cost.',
        'col': AMBER, 'included': False,
    },
]

for i, c in enumerate(costs):
    row = i // 3
    col = i % 3
    cx = Inches(.38) + Inches(4.28) * col
    cy = Inches(1.28) + Inches(2.65) * row
    cw = Inches(4.08)
    ch = Inches(2.38)
    card(sl, cx, cy, cw, ch, accent=c['col'])

    inc_lbl = '✓  INCLUDED IN BUILD' if c['included'] else '⚡  BILLED AT COST'
    inc_col = GREEN if c['included'] else AMBER
    add_text(sl, inc_lbl,
             cx + Inches(.18), cy + Inches(.18), cw - Inches(.3), Inches(.24),
             size=Pt(8), bold=True, color=inc_col, font_name='Courier New')
    add_text(sl, c['label'],
             cx + Inches(.18), cy + Inches(.44), cw - Inches(.3), Inches(.35),
             size=Pt(13), bold=True, color=WHITE, font_name='Calibri')
    add_text(sl, c['cost'],
             cx + Inches(.18), cy + Inches(.82), cw - Inches(.3), Inches(.45),
             size=Pt(20), bold=True, color=c['col'], font_name='Calibri')
    add_text(sl, c['note'],
             cx + Inches(.18), cy + Inches(1.32), cw - Inches(.3), Inches(.9),
             size=Pt(9), color=TEXT_MID, font_name='Calibri')

add_text(sl, '★  InvenIQ does not mark up OpenAI API costs. All AI inference charges pass through directly to your OpenAI account at their published rates.',
         Inches(.38), Inches(6.9), Inches(12.55), Inches(.4),
         size=Pt(9), color=TEXT_DIM, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DEPENDENT COMPONENTS ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, GREEN)

add_text(sl, 'Dependent Components — Architecture Map',
         Inches(.45), Inches(.3), Inches(10), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'How Design Quote Studio connects to the rest of InvenIQ',
         Inches(.45), Inches(.82), Inches(10), Inches(.28),
         size=Pt(9), color=TEXT_MID, font_name='Calibri')

# Centre DQS box
dqs_x, dqs_y, dqs_w, dqs_h = Inches(5.0), Inches(2.9), Inches(3.33), Inches(1.65)
card(sl, dqs_x, dqs_y, dqs_w, dqs_h, accent=GREEN)
add_text(sl, '🏛  Design Quote Studio',
         dqs_x + Inches(.15), dqs_y + Inches(.18), dqs_w - Inches(.3), Inches(.38),
         size=Pt(13), bold=True, color=GREEN, font_name='Calibri')
add_text(sl, 'activeView: designquote\nRole: architect (restricted)',
         dqs_x + Inches(.15), dqs_y + Inches(.6), dqs_w - Inches(.3), Inches(.85),
         size=Pt(8.5), color=TEXT_MID, font_name='Courier New')

deps = [
    (Inches(.35),  Inches(1.1),  Inches(3.4), Inches(1.55), TEAL,   '🗂  Product Catalog',    'Shared SKU database\nHSN + category inference\nBrand / finish tracking'),
    (Inches(9.6),  Inches(1.1),  Inches(3.35), Inches(1.55), BLUE,   '📝  Quotation Builder',  'QB→DQS item sync\nSave scanned items\nLocalStorage 30-min TTL'),
    (Inches(.35),  Inches(4.55), Inches(3.4), Inches(1.75), PURPLE,  '🤖  GPT-4o AI Engine',   'BOQ scanner extraction\nHSN inference\nPhoto item identification'),
    (Inches(9.6),  Inches(4.55), Inches(3.35), Inches(1.75), AMBER,   '💬  WhatsApp Scanner',   'Multi-photo upload\nPer-image AI parse\nCollapsible room sections'),
    (Inches(.35),  Inches(1.1) + Inches(3.72), Inches(3.4), Inches(1.6), RGBColor(0x06,0x78,0x60), '📄  PDF Generator',  'Client-side rendering\nCustom branding (Enterprise)\nNo server-side cost'),
    (Inches(9.6),  Inches(1.1) + Inches(3.72), Inches(3.35), Inches(1.6), RGBColor(0x9c,0x27,0xb0), '📊  MySQL Database',  'sales_design_quotes\nquote_line_items\narchitect_proposals'),
]

for dx, dy, dw, dh, col, title, desc in deps:
    card(sl, dx, dy, dw, dh, accent=col)
    add_text(sl, title, dx + Inches(.15), dy + Inches(.18), dw - Inches(.3), Inches(.35),
             size=Pt(12), bold=True, color=col, font_name='Calibri')
    add_text(sl, desc,  dx + Inches(.15), dy + Inches(.55), dw - Inches(.3), Inches(1.0),
             size=Pt(9), color=TEXT_MID, font_name='Calibri')

# connector lines (simple rectangles)
line_col = BORDER
add_rect(sl, Inches(3.75), Inches(3.3),  Inches(1.25), Pt(1.2), fill_rgb=TEXT_DIM)  # left → centre
add_rect(sl, Inches(8.33), Inches(3.3),  Inches(1.27), Pt(1.2), fill_rgb=TEXT_DIM)  # centre → right
add_rect(sl, Inches(3.75), Inches(5.5),  Inches(1.25), Pt(1.2), fill_rgb=TEXT_DIM)  # bottom-left
add_rect(sl, Inches(8.33), Inches(5.5),  Inches(1.27), Pt(1.2), fill_rgb=TEXT_DIM)  # bottom-right


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 30-DAY FREE TRIAL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
top_bar(sl, GREEN)

add_text(sl, '30-Day Free Trial',
         Inches(.45), Inches(.3), Inches(10), Inches(.5),
         size=Pt(22), bold=True, color=WHITE, font_name='Calibri')
add_text(sl, 'No credit card required · Professional plan features · Activate immediately',
         Inches(.45), Inches(.82), Inches(10), Inches(.28),
         size=Pt(9), color=TEXT_MID, font_name='Calibri')

# big gift box left
add_rect(sl, Inches(.35), Inches(1.2), Inches(4.5), Inches(5.8),
         fill_rgb=RGBColor(0x05, 0x2e, 0x16))
add_rect(sl, Inches(.35), Inches(1.2), Inches(.06), Inches(5.8), fill_rgb=GREEN)
add_text(sl, '🎁', Inches(.6), Inches(1.6), Inches(4.0), Inches(1.5),
         size=Pt(72), color=GREEN, align=PP_ALIGN.CENTER, font_name='Segoe UI Emoji')
add_text(sl, 'Professional Plan\nFeatures Included',
         Inches(.6), Inches(3.2), Inches(4.0), Inches(.9),
         size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(sl, '30 Days · Free · No Card',
         Inches(.6), Inches(4.12), Inches(4.0), Inches(.4),
         size=Pt(13), color=GREEN, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(sl, '→  Unlimited quotations and proposals\n'
             '→  3 architect seat trial\n'
             '→  AI BOQ Scanner unlimited\n'
             '→  QB↔DQS sync active\n'
             '→  Full product catalog access',
         Inches(.6), Inches(4.65), Inches(4.0), Inches(2.1),
         size=Pt(10), color=TEXT_MID, align=PP_ALIGN.LEFT, font_name='Calibri')

# right — activation steps
steps = [
    (GREEN,  '1', 'Architect Role Is Pre-Provisioned',
              'The architect user account (username: architect, password: arch@2026) is\nalready created in every InvenIQ installation. No extra setup needed.'),
    (BLUE,   '2', 'Navigate to Settings → Users',
              'Log in as admin, go to Settings, and modify the architect user\'s password\nand email. Optionally create additional architect accounts.'),
    (PURPLE, '3', 'Log In as Architect',
              'The architect role sees only 3 modules: Design Quote Studio, Settings,\nand About. Full DQS access begins immediately on first login.'),
    (AMBER,  '4', 'Start Creating Quotations',
              'Create your first Interior Quotation. Use the WhatsApp scanner to import\na site photograph BOQ. Export as PDF and send to your client.'),
]
for i, (col, num, title, desc) in enumerate(steps):
    sy = Inches(1.22) + Inches(1.5) * i
    card(sl, Inches(5.2), sy, Inches(7.85), Inches(1.32), accent=col)
    add_rect(sl, Inches(5.2), sy + Inches(.35), Inches(.38), Inches(.38),
             fill_rgb=col)
    add_text(sl, num, Inches(5.2), sy + Inches(.35), Inches(.38), Inches(.38),
             size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Calibri')
    add_text(sl, title,
             Inches(5.68), sy + Inches(.18), Inches(7.2), Inches(.35),
             size=Pt(12), bold=True, color=col, font_name='Calibri')
    add_text(sl, desc,
             Inches(5.68), sy + Inches(.55), Inches(7.2), Inches(.72),
             size=Pt(9.5), color=TEXT_MID, font_name='Calibri')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONTACT & NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_DARK)
add_rect(sl, 0, 0, SLIDE_W, Inches(.07), fill_rgb=GREEN)

# full-width green band at bottom
add_rect(sl, 0, Inches(5.8), SLIDE_W, Inches(1.7), fill_rgb=RGBColor(0x05, 0x2e, 0x16))
add_rect(sl, 0, Inches(7.43), SLIDE_W, Pt(3), fill_rgb=GREEN)

add_text(sl, 'Ready to Get Started?',
         Inches(.6), Inches(.5), Inches(12.1), Inches(.65),
         size=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(sl, 'Design Quote Studio is live in your InvenIQ build today',
         Inches(.6), Inches(1.18), Inches(12.1), Inches(.38),
         size=Pt(13), color=TEXT_MID, align=PP_ALIGN.CENTER, font_name='Calibri')

next_steps = [
    (TEAL,   '🚀', 'Activate Free Trial',     'Log in as architect user\nand start your 30-day\nProfessional trial today'),
    (BLUE,   '📞', 'Book a Live Demo',         'Schedule a 30-minute DQS\ndemo with the InvenIQ\nproduct team'),
    (PURPLE, '📄', 'Download Full Brochure',   'Get the complete DQS\nproduct brochure with\nclient case studies'),
    (GREEN,  '⬆',  'Upgrade Your Plan',        'Contact your account\nmanager to upgrade from\nStarter to Professional'),
]

for i, (col, icon, title, desc) in enumerate(next_steps):
    nx = Inches(.38) + Inches(3.25) * i
    card(sl, nx, Inches(1.82), Inches(3.06), Inches(3.68), accent=col)
    add_text(sl, icon, nx + Inches(.15), Inches(2.05), Inches(2.75), Inches(.75),
             size=Pt(36), color=col, font_name='Segoe UI Emoji')
    add_text(sl, title, nx + Inches(.15), Inches(2.82), Inches(2.75), Inches(.45),
             size=Pt(13), bold=True, color=col, font_name='Calibri')
    add_text(sl, desc, nx + Inches(.15), Inches(3.3), Inches(2.75), Inches(1.1),
             size=Pt(10), color=TEXT_MID, font_name='Calibri')

# bottom bar contact info
contact = [
    ('Product', 'InvenIQ v3.7'),
    ('Module',  'Design Quote Studio'),
    ('Role',    'Architect (restricted)'),
    ('AI',      'GPT-4o · OpenAI'),
    ('Stack',   'React 18 + FastAPI + MySQL'),
    ('Date',    'June 2026'),
]
for i, (k, v) in enumerate(contact):
    bx = Inches(.5) + Inches(2.22) * i
    add_text(sl, k.upper(), bx, Inches(5.98), Inches(2.0), Inches(.24),
             size=Pt(7.5), bold=True, color=TEXT_DIM, font_name='Courier New')
    add_text(sl, v, bx, Inches(6.24), Inches(2.0), Inches(.28),
             size=Pt(9.5), color=WHITE, font_name='Calibri')

add_text(sl, 'InvenIQ — Design Quote Studio · Subscription & Licensing Overview',
         Inches(.38), Inches(5.96), Inches(12.55), Inches(.28),
         size=Pt(9), color=TEXT_DIM, align=PP_ALIGN.CENTER, font_name='Calibri')


# ── Save ───────────────────────────────────────────────────────────────────
out = r'c:\InvenIQ\DQS_Subscription_Pricing.pptx'
prs.save(out)
print(f'Saved: {out}')
